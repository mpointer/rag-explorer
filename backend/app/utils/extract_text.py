from __future__ import annotations
import io
import json
import csv
from typing import Optional

from bs4 import BeautifulSoup
from pypdf import PdfReader
from docx import Document
from pptx import Presentation


def _decode_text(data: bytes) -> str:
    try:
        return data.decode("utf-8")
    except Exception:
        return data.decode("latin-1", errors="ignore")


def extract_text(file_bytes: bytes, filename: Optional[str] = None, content_type: Optional[str] = None) -> str:
    name = (filename or '').lower()
    ct = (content_type or '').lower()

    # Plain text
    if name.endswith(('.txt', '.md', '.csv', '.json')) or ct.startswith('text/'):
        if name.endswith('.csv'):
            sio = io.StringIO(_decode_text(file_bytes))
            reader = csv.reader(sio)
            return "\n".join([", ".join(row) for row in reader])
        if name.endswith('.json'):
            try:
                obj = json.loads(_decode_text(file_bytes))
                return json.dumps(obj, indent=2)
            except Exception:
                return _decode_text(file_bytes)
        return _decode_text(file_bytes)

    # PDF
    if name.endswith('.pdf') or ct in ('application/pdf',):
        reader = PdfReader(io.BytesIO(file_bytes))
        texts = []
        for page in reader.pages:
            try:
                texts.append(page.extract_text() or "")
            except Exception:
                continue
        return "\n\n".join(texts).strip()

    # DOCX
    if name.endswith('.docx') or ct in ('application/vnd.openxmlformats-officedocument.wordprocessingml.document',):
        doc = Document(io.BytesIO(file_bytes))
        paras = [p.text for p in doc.paragraphs if p.text]
        return "\n".join(paras).strip()

    # PPTX
    if name.endswith('.pptx') or ct in ('application/vnd.openxmlformats-officedocument.presentationml.presentation',):
        prs = Presentation(io.BytesIO(file_bytes))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    t = getattr(shape, "text", "")
                    if t:
                        texts.append(t)
        return "\n\n".join(texts).strip()

    # HTML
    if name.endswith(('.html', '.htm')) or ct in ('text/html',):
        soup = BeautifulSoup(_decode_text(file_bytes), 'html.parser')
        return soup.get_text(" ", strip=True)

    return _decode_text(file_bytes)
